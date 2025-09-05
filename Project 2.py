import pandas as pd
import numpy as np
import os
import matplotlib.pyplot as plt
import seaborn as sns
import plotly.express as px
import plotly.io as pio
from pptx import Presentation
from pptx.util import Inches, Pt

# ------------------ CONFIG ------------------
INPUT_FILE = "C:/Users/Abdullah Umer/Desktop/CodeAlpha Internship/Project 2/Global Population Dataset.csv"
OUTPUT_DIR = "C:/Users/Abdullah Umer/Desktop/CodeAlpha Internship/Project 2/outputs"
os.makedirs(OUTPUT_DIR, exist_ok=True)

# ------------------ LOAD + CLEAN ------------------
df = pd.read_csv(INPUT_FILE)
df.columns = [c.strip() for c in df.columns]

# Detect population columns
pop_cols = [c for c in df.columns if c.lower().startswith("pop") and any(ch.isdigit() for ch in c)]
id_vars = [c for c in df.columns if c not in pop_cols]

# Reshape
long_df = df.melt(id_vars=id_vars, value_vars=pop_cols, var_name="year_raw", value_name="population")
long_df["year"] = long_df["year_raw"].astype(str).str.extract(r"(\d{4})").astype(int)
long_df = long_df.dropna(subset=["population"])

# Pivot for convenience
pivot = long_df.pivot_table(index="country", columns="year", values="population", aggfunc="first")

# Pick latest observed year ≤ 2024
latest_year = max([y for y in pivot.columns if y <= 2024]) if any(y <= 2024 for y in pivot.columns) else max(pivot.columns)

# ------------------ BASIC METRICS ------------------
top5 = pivot[latest_year].nlargest(5)

if 2010 in pivot.columns and 2024 in pivot.columns:
    span = 2024 - 2010
    cagr_df = ((pivot[2024] / pivot[2010]) ** (1 / span) - 1) * 100
    top_cagr = cagr_df.nlargest(5)
else:
    cagr_df, top_cagr = None, None

# ------------------ VISUALIZATIONS ------------------
# 1. Line chart (top 7)
top7 = pivot[latest_year].nlargest(7).index
plt.figure(figsize=(10,6))
for c in top7:
    ys = pivot.loc[c].dropna()
    plt.plot(ys.index, ys.values, marker="o", label=c)
plt.title(f"Population over time — Top 7 countries in {latest_year}")
plt.xlabel("Year"); plt.ylabel("Population"); plt.legend(); plt.grid(True)
plt.tight_layout()
line_path = os.path.join(OUTPUT_DIR, "line_top7.png")
plt.savefig(line_path, dpi=150); plt.close()

# 2. Bar chart (top 15)
top15 = pivot[latest_year].nlargest(15)
plt.figure(figsize=(10,6))
plt.bar(top15.index, top15.values)
plt.xticks(rotation=45, ha="right")
plt.title(f"Top 15 countries by population in {latest_year}")
plt.ylabel("Population")
plt.tight_layout()
bar_path = os.path.join(OUTPUT_DIR, "bar_top15.png")
plt.savefig(bar_path, dpi=150); plt.close()

# 3. Histogram
plt.figure(figsize=(8,5))
plt.hist(pivot[latest_year].dropna(), bins=30)
plt.title(f"Histogram of populations in {latest_year}")
plt.xlabel("Population"); plt.ylabel("# of countries")
plt.tight_layout()
hist_path = os.path.join(OUTPUT_DIR, "hist_population.png")
plt.savefig(hist_path, dpi=150); plt.close()

# 4. Pie chart (top 10 + Others)
top10 = df.nlargest(10, f"pop{latest_year}").set_index("country")
others = max(df[f"pop{latest_year}"].sum() - top10[f"pop{latest_year}"].sum(), 0)
vals = list(top10[f"pop{latest_year}"]) + [others]
labs = list(top10.index) + ["Others"]
plt.figure(figsize=(8,8))
plt.pie(vals, labels=labs, autopct="%1.1f%%", startangle=140)
plt.title("Share of World Population — Top 10 + Others")
plt.tight_layout()
pie_path = os.path.join(OUTPUT_DIR, "pie_top10_share.png")
plt.savefig(pie_path, dpi=150); plt.close()

# 5. Scatter population vs density (if density available)
scatter_path = os.path.join(OUTPUT_DIR, "scatter_pop_density.png")
if "density" in df.columns and f"pop{latest_year}" in df.columns:
    plt.figure(figsize=(8,6))
    plt.scatter(df[f"pop{latest_year}"], df["density"], alpha=0.6)
    plt.xscale("log")
    plt.xlabel(f"Population {latest_year} [log]"); plt.ylabel("Density")
    plt.title(f"Population vs Density — {latest_year}")
    plt.tight_layout()
    plt.savefig(scatter_path, dpi=150)
    plt.close()
else:
    scatter_path = None

# 6. Correlation heatmap
num_cols = df.select_dtypes(include=[np.number]).columns
df_corr = df[num_cols].corr()
plt.figure(figsize=(10,8))
sns.heatmap(df_corr, cmap="coolwarm", annot=False)
plt.title("Correlation heatmap (numeric columns)")
plt.tight_layout()
heatmap_path = os.path.join(OUTPUT_DIR, "heatmap_corr.png")
plt.savefig(heatmap_path, dpi=150); plt.close()

# 7. Interactive Choropleth (Plotly)
if "cca3" in df.columns and f"pop{latest_year}" in df.columns:
    fig = px.choropleth(df, locations="cca3", color=f"pop{latest_year}", hover_name="country",
                        title=f"World population by country ({latest_year})")
    pio.write_html(fig, file=os.path.join(OUTPUT_DIR, f"choropleth_{latest_year}.html"), auto_open=False)

# 8. Interactive line chart top 10
int_df = long_df[long_df["country"].isin(pivot[latest_year].nlargest(10).index)]
fig_line = px.line(int_df, x="year", y="population", color="country", markers=True,
                   title=f"Population trends — Top 10 countries ({latest_year})")
pio.write_html(fig_line, file=os.path.join(OUTPUT_DIR, "interactive_line_top10.html"), auto_open=False)

# ------------------ REPORTING ------------------
report = []
report.append(f"Years available: {sorted(long_df['year'].unique())}")
report.append(f"Top 5 countries by population in {latest_year}:")
for i,(c,v) in enumerate(top5.items(),1):
    report.append(f"{i}. {c}: {int(v):,}")

if top_cagr is not None:
    report.append("\nTop 5 by CAGR (2010–2024):")
    for i,(c,v) in enumerate(top_cagr.items(),1):
        report.append(f"{i}. {c}: {v:.2f}%/yr")

with open(os.path.join(OUTPUT_DIR, "population_insights.md"), "w") as f:
    f.write("\n".join(report))

# ------------------ POWERPOINT ------------------
prs = Presentation()
# Title slide
title_slide_layout = prs.slide_layouts[0]
slide = prs.slides.add_slide(title_slide_layout)
title, subtitle = slide.shapes.title, slide.placeholders[1]
title.text = "Global Population Data Visualization"
subtitle.text = "Auto-generated report"

# Add slides for images
img_files = [line_path, bar_path, hist_path, pie_path, scatter_path, heatmap_path]
for img in img_files:
    if img is not None and os.path.exists(img):
        blank_slide_layout = prs.slide_layouts[6]  # blank slide
        slide = prs.slides.add_slide(blank_slide_layout)
        slide.shapes.add_picture(img, Inches(1), Inches(1.5), width=Inches(8))
    elif img is not None:
        print(f"⚠️ Image not found: {img}")

# Save PPTX
pptx_path = os.path.join(OUTPUT_DIR, "population_report.pptx")
prs.save(pptx_path)
print(f"✅ PowerPoint saved at: {pptx_path}")











